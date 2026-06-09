#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
文件解析工具

支持格式：
- PDF (.pdf) - 使用 PyMuPDF
- Word (.docx) - 使用 docx2txt
- Word (.doc) - 使用 catdoc（沙箱需预装 catdoc）
- Excel (.xlsx) - 使用 openpyxl
- Excel (.xls) - 使用 xlrd
- PowerPoint (.pptx) - 使用 python-pptx

用法：
    python3 file_parser.py parse <file_path>
"""

import sys
import os
import json


def parse_pdf(file_path):
    """解析 PDF 文件"""
    import fitz
    doc = fitz.open(file_path)
    text_parts = []
    for page_num, page in enumerate(doc, 1):
        text = page.get_text()
        if text.strip():
            text_parts.append(f"--- 第 {page_num} 页 ---\n{text}")
    doc.close()
    return "\n\n".join(text_parts)


def parse_docx(file_path):
    """解析 Word (.docx) 文件"""
    import docx2txt
    text = docx2txt.process(file_path)
    return text


def parse_doc(file_path):
    """解析旧版 Word (.doc) 文件 — 通过 catdoc 命令行工具

    catdoc 是专门的 .doc 解析工具（OLE2 复合文档），体积小（~1MB），
    中文支持好（通过 -s cp936 -d utf-8 指定编码转换）。

    当前沙箱运行时不安装系统依赖；如需支持 .doc，请在镜像中预置 catdoc。
    """
    import subprocess
    import shutil

    # 检查 catdoc 是否安装
    catdoc_path = shutil.which("catdoc")
    if catdoc_path is None:
        raise RuntimeError(
            "catdoc 未安装，无法解析旧版 .doc 文件。请在沙箱镜像中预置 catdoc，"
            "或将文件转换为 .docx/.pdf 后再解析。"
        )

    result = subprocess.run(
        [catdoc_path, "-s", "cp936", "-d", "utf-8", file_path],
        capture_output=True,
        text=True,
        timeout=60
    )

    if result.returncode != 0:
        raise RuntimeError(f"catdoc 解析失败: {result.stderr or 'unknown error'}")

    return result.stdout


def parse_xlsx(file_path):
    """解析 Excel (.xlsx) 文件"""
    from openpyxl import load_workbook
    wb = load_workbook(file_path, data_only=True)
    result_parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        result_parts.append(f"--- Sheet: {sheet_name} ---")
        for row in ws.iter_rows(values_only=True):
            row_str = "\t".join([str(cell) if cell is not None else "" for cell in row])
            if row_str.strip():
                result_parts.append(row_str)
    wb.close()
    return "\n".join(result_parts)


def parse_xls(file_path):
    """解析 Excel (.xls) 文件"""
    import xlrd
    wb = xlrd.open_workbook(file_path)
    result_parts = []
    for sheet in wb.sheets():
        result_parts.append(f"--- Sheet: {sheet.name} ---")
        for row_idx in range(sheet.nrows):
            row = [str(sheet.cell_value(row_idx, col_idx)) for col_idx in range(sheet.ncols)]
            row_str = "\t".join(row)
            if row_str.strip():
                result_parts.append(row_str)
    return "\n".join(result_parts)


def parse_pptx(file_path):
    """解析 PowerPoint (.pptx) 文件"""
    from pptx import Presentation
    prs = Presentation(file_path)
    result_parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        result_parts.append(f"--- 幻灯片 {slide_num} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        result_parts.append(text)
    return "\n".join(result_parts)


def parse_file(file_path):
    """根据文件扩展名选择解析方法"""
    if not os.path.exists(file_path):
        return json.dumps({"success": False, "error": f"文件不存在: {file_path}"}, ensure_ascii=False)

    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext == ".pdf":
            text = parse_pdf(file_path)
        elif ext == ".docx":
            text = parse_docx(file_path)
        elif ext == ".doc":
            text = parse_doc(file_path)
        elif ext == ".xlsx":
            text = parse_xlsx(file_path)
        elif ext == ".xls":
            text = parse_xls(file_path)
        elif ext == ".pptx":
            text = parse_pptx(file_path)
        else:
            return json.dumps({"success": False, "error": f"不支持的文件格式: {ext}"}, ensure_ascii=False)

        return json.dumps({"success": True, "format": ext, "content": text}, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"success": False, "error": str(e)}, ensure_ascii=False)


def main():
    if len(sys.argv) < 3:
        print("用法: python3 file_parser.py parse <file_path>")
        sys.exit(1)

    command = sys.argv[1]

    if command == "parse":
        file_path = sys.argv[2]
        result = parse_file(file_path)
        print(result)
    else:
        print(f"未知命令: {command}")
        sys.exit(1)


if __name__ == "__main__":
    main()
